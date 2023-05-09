import os
from typing import Dict, Union
from pptx import Presentation
from pptx.shapes.base import BaseShape


class PptxParser:
    def __init__(self, pptx_file_path: str):
        if not os.path.exists(pptx_file_path):
            raise FileNotFoundError(f"File not found: {pptx_file_path}")
        self.pptx_file_path = pptx_file_path

    def _get_text_from_shape(self, shape: BaseShape) -> str:
        text = ""
        if shape.has_text_frame:
            paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs]
            text = " ".join(paragraphs)
        return text

    def parse_presentation(self) -> Dict[int, str]:
        ppt = Presentation(self.pptx_file_path)
        slides_dict = {}

        for slide_num, slide in enumerate(ppt.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                text = self._get_text_from_shape(shape)
                if text:
                    slide_text.append(text)
            slides_dict[slide_num] = " ".join(slide_text)

        return slides_dict



def main():

    pptx_file_path = "./test/files/presentation_for_tst.pptx"

    # Create an instance of PptxParser with the file path
    pptx_parser = PptxParser(pptx_file_path)

    # Parse the presentation and store the resulting dictionary
    slides_dict = pptx_parser.parse_presentation()

    # Print the dictionary
    for slide_num, slide_text in slides_dict.items():
        print(f"Slide {slide_num}:")
        print(slide_text)
        print()

if __name__ == "__main__":
    main()
