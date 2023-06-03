from typing import Dict
from pptx import Presentation
from pptx.shapes.base import BaseShape


class PptxParser:
    def __init__(self, pptx_file_path: str):
        self.ppt = Presentation(pptx_file_path)

    def _get_text_from_shape(self, shape: BaseShape) -> str:
        if shape.has_text_frame:
            return " ".join(p.text.strip() for p in shape.text_frame.paragraphs)
        return ""

    def parse_presentation(self) -> Dict[int, str]:
        return {
            slide_num: " ".join(
                self._get_text_from_shape(shape)
                for shape in slide.shapes
                if self._get_text_from_shape(shape)
            )
            for slide_num, slide in enumerate(self.ppt.slides, start=1)
        }


def main():
    pptx_file_path = "../Test/files/presentation_for_tst.pptx"
    pptx_parser = PptxParser(pptx_file_path)
    slides_dict = pptx_parser.parse_presentation()

    for slide_num, slide_text in slides_dict.items():
        print(f"Slide {slide_num}:\n{slide_text}\n")


if __name__ == "__main__":
    main()
